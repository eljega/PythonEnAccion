import os
from pptx import Presentation
import fitz
from docx import Document

"""
Extractor de texto e imagenes

este script extrae texto imagenes y tablas de documentos PDF, word y powerpoint
y los organiza en carpetas separadas segun el tipo de contenido.

Dependencias a instalar:
- pip install python-docx PyMuPDF python-pptx

Instrucciones de Uso:
1. Coloca los documentos de los que deseas extraer contenido en un directorio accesible.
2. Asegurate de modificar las rutas de los archivos en las variables.
"""

def extraer_de_pdf(ruta_pdf):
    doc = fitz.open(ruta_pdf)
    texto = ""
    for pagina in doc:
        texto += pagina.get_text()
    
    #para guardar texto extraido de PDF
    with open(os.path.join("ruta para guardar", "pdf_texto.txt"), "w", encoding='utf-8') as f:
        f.write(texto)
    
    for i, pagina in enumerate(doc):
        imagenes = pagina.get_images(full=True)
        for img_index, img in enumerate(imagenes, start=1):
            base_image = doc.extract_image(img[0])
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_path = f"C:escribe aqui la ruta para guardar{i+1}_img_{img_index}.{image_ext}"
            with open(image_path, "wb") as f:
                f.write(image_bytes)

def extraer_de_word(ruta_docx):
    doc = Document(ruta_docx)
    texto = "\n".join([p.text for p in doc.paragraphs])
    with open(os.path.join("ruta para guardar", "word_texto.txt"), "w", encoding='utf-8') as f:
        f.write(texto)
    
def extraer_de_ppt(ruta_pptx):
    prs = Presentation(ruta_pptx)
    texto = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    
    with open(os.path.join("ruta para guardar", "ppt_texto.txt"), "w", encoding='utf-8') as f:
        f.write(texto)
    

if __name__ == "__main__":
    #recuerda modificar las rutas segun tus carpetas
    ruta_pdf = "ruta donde tienes tu .pdf"
    ruta_docx = "ruta donde tienes tu .docx"
    ruta_pptx = "ruta donde tienes tu .pptx"
    
    extraer_de_pdf(ruta_pdf)
    extraer_de_word(ruta_docx)
    extraer_de_ppt(ruta_pptx)
